import os
import re
from enum import Enum
from functools import cached_property
from pathlib import Path
from typing import Iterable, TypedDict

if __name__ == "__main__":
    if Path(os.getcwd()).name == "processing":
        os.chdir("../..")

from backend.processing.checker import Checker
from pptx import Presentation
from pptx.slide import Slide
from thefuzz import fuzz


class Status(str, Enum):
    ERROR = "Error"
    WARNING = "Warning"
    PASS = "Passing"

    def __repr__(self):
        if self == Status.ERROR:
            return "Error"
        elif self == Status.WARNING:
            return "Warning"
        elif self == Status.PASS:
            return "Pass"


class Result(TypedDict):
    title: str
    comments: str
    status: Status


SlideSubset = dict[int, Slide]

section_mapping = {
    "Bible Reading": "Hearing God\u2019s Word Read",
    "Sermon": "Hearing God\u2019s Word Proclaimed",
    "Discuss in groups": "Sermon Discussion",
}


def raw_req_order_of_service_no_declaration() -> str:
    return """Opening Words	1	
Opening Song	4	Behold Our God
Family Confession	2	#11 Confession of Sin (Slide 17 & 18)
Family Prayer	4	Refer to Prayer Points Tab in this document (Usually updated by Thu)
Family Business	5	Refer to Family Business Tab
Bible Reading 	4	Daniel 5
Sermon	30	Preacher: Denesh
Closing Song	4	Only a Holy God
Closing Words	1	
Discuss in groups	5	
Dismissal		"""


def get_clean_req_order_of_service(
    raw_req_order_of_service_no_declaration: str,
) -> list[tuple[str, int, str]]:
    """
    Returns the cleaned up required order of service.

    Args:
        raw_req_order_of_service_no_declaration (str): Raw required order of service

    Returns:
        list[tuple[str, int, str]]: _description_
    """
    result, intermediate = [], []
    for item in raw_req_order_of_service_no_declaration.split("\n"):
        intermediate.append(item.split("\t"))
    for item in intermediate:
        result.append((item[0].strip(), item[1].strip(), item[2].strip()))
    return result


def filter_clean_req_order_of_service(
    clean_req_order_of_service: list[tuple[str, int, str]],
) -> list[tuple[str, str]]:
    """
    Returns the filtered clean required order of service with elements at position:
    1. Item
    2. Things to note

    Returns:
        tuple[str, str]

    Args:
        raw_req_order_of_service_no_declaration (str): Raw

    Returns:
        list[tuple[str, str]]: _description_
    """
    return [
        (section_mapping.get(item[0], item[0]), item[2])
        for item in clean_req_order_of_service
        if not re.match("(Opening Words|Dismissal|Closing Words)", item[0])
    ]


class ContentChecker(Checker):
    """
    Checks the content of the uploaded slides according to the inputs.
    """

    def __init__(
        self,
        selected_date: str,
        req_order_of_service: str,
        sermon_discussion_qns: str,
        presentations: list,
    ) -> None:
        self.selected_date = selected_date
        self.req_order_of_service = req_order_of_service
        self.sermon_discussion_qns = sermon_discussion_qns
        self.presentations = presentations

    @cached_property
    def pptx(self):
        # TODO: extend to multiple presentations
        return self.presentations[0]

    def run(self) -> list[Result]:
        """
        Runs all the checks within the ContentChecker.

        Returns:
            list[dict[str,str]]: List of dictionaries containing 'title' and 'description' keys
        """
        clean_req_order_of_service = filter_clean_req_order_of_service(
            get_clean_req_order_of_service(self.req_order_of_service)
        )
        result = [
            self.check_existence_of_section_headers(
                section_headers=section_headers(self.pptx.slides)
            ),
            *self.check_section_headers_have_correct_order(
                req_order_of_service=clean_req_order_of_service,
                slide_order_of_service=slide_order_of_service(
                    section_headers(self.pptx.slides)
                ),
            ),
            *self.check_all_dates_are_as_provided(date=self.selected_date),
        ]
        return result

    def check_existence_of_section_headers(
        self, section_headers: SlideSubset
    ) -> Result:
        result: Result = {
            "title": "Check existence of section header slides",
            "comments": f"{len(section_headers)} section header slides found.",
            "status": Status.PASS if len(section_headers) > 0 else Status.ERROR,
        }
        return result

    def check_section_headers_have_correct_order(
        self,
        req_order_of_service: list[tuple[str, str]],
        slide_order_of_service: dict[int, list[str]],
    ) -> list[Result]:
        """
        Test all section headers have the correct order of service by checking that the
        order of service provided in each section header slide matches the provided order
        of service.

        Section headers are the slides that separate each section of the service.
        These slides will usually have the order of service in grey on the right side.

        Args:
            req_order_of_service (tuple[str, str]): Required order of service
            slide_order_of_service (list[tuple[str, str]]): Order of services from slides
        """
        # 1. Identify the section header slides
        # 2. Extract the text from the slide
        # 3. Check that the order of service from the text in the slides is correct
        results: list[Result] = []
        items_with_comments = (
            "(Opening Song|Closing Song|Hearing God(\u2018|\u2019|')s Word Read)"
        )

        for i, slide_text in slide_order_of_service.items():
            index = 0
            for entry in slide_text:

                title, comments = req_order_of_service[index]
                title = title.replace("\u2018", "\u2019")
                required_item = f"{title} \u2013 {comments}"
                is_commented_item = re.match(items_with_comments, entry)
                partial_ratio = fuzz.partial_ratio(required_item, entry)

                if is_commented_item is not None:
                    is_commented_items_correct = required_item == entry
                    if is_commented_items_correct:
                        index += 1
                        continue
                    elif 90 < partial_ratio < 100:
                        result = {
                            "title": "Is there a typo?",
                            "status": Status.WARNING,
                            "comments": f"On Slide {i}, Expected: '{required_item}'. Provided: '{entry}'. Similarity score = {partial_ratio} of 100",
                        }
                        index += 1
                    else:
                        result = {
                            "title": "Check section headers are in the correct order",
                            "comments": f"On Slide {i}, Expected: '{required_item}'. Provided: '{entry}'. Similarity score = {partial_ratio} of 100",
                            "status": Status.ERROR,
                        }
                    if result not in results:
                        results.append(result)
                elif "\u2018" in entry:
                    result = {
                        "title": "Is there a typo?",
                        "comments": f"On Slide {i}, Expected: '{required_item}'. Provided: '{entry}'. The use of the unicode character U+2018 (\u2018) is triggering this warning; replace this character with U+2019 (\u2019) or a standard single quote (') to resolve this error. Similarity score = {partial_ratio} of 100",
                        "status": Status.WARNING,
                    }
                    results.append(result)
                    index += 1
                elif entry.replace("\u2018", "\u2019") != title:
                    continue
                else:
                    index += 1

        if len(results) == 0:
            result: Result = {
                "title": "Check all required order of service items are present and in the correct order",
                "status": Status.PASS,
                "comments": "All slides containing order of service have the required order of service items and are presented in the correct order.",
            }
            results.append(result)
        return results

    def check_family_confession_content_matches_number(
        self, req_order_of_service: tuple[str, str]
    ):
        """
        Test family confession has the correct contents by checking that the words on the
        slide match the required content (usually a #number) specified in the order of
        service.

        Family Confession is an item in the order of service.
        """
        # 1. Identify the family confession slides
        # 2. Check that there is a text box on the page that contains the family confession
        # 3. Check that it matches the number
        pass

    def check_family_declaration_content_matches_number(
        self,
        req_order_of_service: tuple[str, str],
    ):
        """
        Test that family declaration (if present) has the correct contents by checking that
        the words on the slide match the required content (usually a #number) specified in
        the order of service.

        Family Declaration is an optional item in the order of service.
        """
        pass

    def check_all_lyric_slides_have_no_title(
        self, req_order_of_service: tuple[str, str]
    ):
        """
        Test all lyric slides do not contain a title.

        Lyric slides are slides for the opening and closing song.
        """
        # 1. Identify the lyric slides for the opening and closing song
        # 2. Check that the title is not present (i.e. none of the text boxes contain the song title exclusively)
        pass

    def check_all_dates_are_as_provided(self, date: str):
        """
        Test all dates that appear in the presentation are equal to the date of Sunday service.

        Searches for certain date patterns in the slides.

        Date patterns:
            "01-Jan-2022"
            "01 Jan 2022"

        Args:
            date (str): Date of Sunday service
        """
        date_pattern = "\\d+[\\s-][A-Za-z]+[\\s-]\\d+"
        slides_with_dates = get_slides_by_pattern(self.pptx.slides, date_pattern)
        results = []
        for i, item_list in get_raw_text_from_slides(slides_with_dates).items():
            for item in item_list:
                partial_ratio = fuzz.partial_ratio(item, date)
                if re.match(date_pattern, item) and item.replace(
                    "_", " "
                ) != date.replace("_", " "):
                    result = {
                        "title": "Check all dates that appear in the slides are the same as the date of Sunday service.",
                        "status": Status.ERROR,
                        "comments": f"On slide {i}, Expected: '{date}'. Provided: '{item}'. Similarity score = {partial_ratio} of 100",
                    }
                    results.append(result)

        if len(results) == 0:
            result = {
                "title": "Check all dates that appear in the slides are the same as the date of Sunday service.",
                "status": Status.PASS,
                "comments": "All slides containing dates display the required date of Sunday service.",
            }
            results.append(result)

        return results


def get_slides_by_pattern(all_slides: Iterable[Slide], pattern: str) -> SlideSubset:
    """
    Returns a subset of all slides that contain the provided text argument on the slide

    Args:
        all_slides (Iterable[Slide]): An iterable containing all slides
        pattern (str): String to match

    Returns:
        SlideSubset: Subset of slides with slide number (1-indexed) as keys
    """
    subset = dict()
    for i, slide in enumerate(all_slides, 1):
        shapes = [*slide.shapes, *slide.slide_layout.shapes]  # type: ignore
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            if pattern in shape.text_frame.text or re.match(
                pattern, shape.text_frame.text
            ):
                subset[i] = slide
    return subset


def section_headers(all_slides: Iterable[Slide]) -> SlideSubset:
    """
    Section header slides are slides that contain the text "Today's order of service".
    This includes the Welcome slide, and all other item slides (i.e. Family Confession,
    Family Prayer etc.)

    Args:
        all_slides (Iterable[Slide]): An iterable containing all slides

    Returns:
        list[Slide]: List of slides matching the subset
    """
    text = "order of service"
    return get_slides_by_pattern(all_slides, text)


def get_raw_text_from_slides(slides: SlideSubset) -> dict[int, list[str]]:
    """
    Returns the raw text from any shapes (including text boxes) in the provided slides.

    Args:
        slides (list[Slide]): List of slides

    Returns:
        dict[int, list[str]]: Raw string extracts according to slide number
    """
    result = {}
    for i, slide in slides.items():
        for shape in [*slide.shapes, *slide.slide_layout.shapes]:  # type: ignore
            if i in result:
                result[i].append(shape.text_frame.text)
            else:
                result[i] = [shape.text_frame.text]
    return result


def slide_order_of_service(section_headers: SlideSubset) -> dict[int, list[str]]:
    def filter_order_of_service_only(
        texts: dict[int, list[str]]
    ) -> dict[int, list[str]]:
        result = {}
        for i, item_list in texts.items():
            for item in item_list:
                if "order of service" in item:
                    result[i].append(item)
                else:
                    result[i] = [item]
        return result

    def split_and_strip(text: str) -> list[str]:
        """
        Splits up raw text on newline characters and strips on both sides of the
        resulting string.
        """
        split_text = re.split("\n+", text)
        return [
            item.strip()
            for item in split_text
            if len(item) and "order of service" not in item
        ]

    section_header_text = get_raw_text_from_slides(section_headers)
    orders_of_service = filter_order_of_service_only(section_header_text)
    return {
        i: split_and_strip(item)
        for i, item_list in orders_of_service.items()
        for item in item_list
    }


if __name__ == "__main__":
    # with open("input/01.05 (9am) service slides.pptx", "rb") as f:
    with open("input/22.05 (10.30am) service slides.pptx", "rb") as f:
        pptx = Presentation(f)

    cc = ContentChecker(
        selected_date="22 May 2022",
        req_order_of_service=raw_req_order_of_service_no_declaration(),
        sermon_discussion_qns="",
        presentations=[pptx],
    )
    results = cc.run()

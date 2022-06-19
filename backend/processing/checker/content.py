import os
import re
from functools import cached_property
from pathlib import Path
from typing import Iterable

if __name__ == "__main__":
    if Path(os.getcwd()).parent.name == "processing":
        os.chdir("../../..")

from backend.processing.checker.base import BaseChecker, BaseMultiChecker
from backend.processing.result import FileResults, Result, Status
from pptx import Presentation as PresentationConstructor
from pptx.presentation import Presentation
from pptx.slide import Slide
from thefuzz import fuzz

section_mapping = {
    "Bible Reading": "Hearing God\u2019s Word Read",
    "Sermon": "Hearing God\u2019s Word Proclaimed",
    "Discuss in groups": "Sermon Discussion",
}

CleanOrderOfService = list[tuple[str, int, str]]
FilteredCleanOrderOfService = list[tuple[str, str]]
SlideOrderOfService = dict[int, list[str]]
SlideSubset = dict[int, Slide]


def raw_req_order_of_service_no_declaration() -> str:
    return """Opening Words	1	
Opening Song	4	Behold Our God
Family Confession	2	#11 Confession of Sin (Slide 17 & 18)
Family Prayer	4	Refer to Prayer Points Tab in this document (Usually updated by Thu)
Family Business	5	Refer to Family Business Tab
Bible Reading 	4	Daniel 5
Sermon	30	Preacher: Denesh
Closing Song	4	Only a Holy God
Closing Words	1	
Discuss in groups	5	
Dismissal		"""


def get_clean_req_order_of_service(
    raw_req_order_of_service_no_declaration: str,
) -> CleanOrderOfService:
    """
    Returns the cleaned up required order of service.

    Args:
        raw_req_order_of_service_no_declaration (str): Raw required order of service

    Returns:
        list[tuple[str, int, str]]: Cleaned up required order of service
    """
    result, intermediate = [], []
    for item in raw_req_order_of_service_no_declaration.split("\n"):
        intermediate.append(item.split("\t"))
    for item in intermediate:
        result.append((item[0].strip(), item[1].strip(), item[2].strip()))
    return result


def filter_clean_req_order_of_service(
    clean_req_order_of_service: CleanOrderOfService,
) -> FilteredCleanOrderOfService:
    """
    Returns the filtered clean required order of service with elements at position:
    1. Item
    2. Things to note

    Returns:
        tuple[str, str]

    Args:
        raw_req_order_of_service_no_declaration (str): Raw order of service

    Returns:
        list[tuple[str, str]]: Filtered clean required order of service
    """
    return [
        (section_mapping.get(item[0], item[0]), item[2])
        for item in clean_req_order_of_service
        if not re.match("(Opening Words|Dismissal|Closing Words)", item[0])
    ]


def get_clean_sermon_discussion_qns(sermon_discussion_qns: str) -> list[str]:
    split_text = re.split(r"\d+\.", sermon_discussion_qns)
    return [text.strip() for text in split_text if text.strip()]


def get_slides_by_pattern(all_slides: Iterable[Slide], pattern: str) -> SlideSubset:
    """
    Returns a subset of all slides that contain the provided text argument on the slide.

    Args:
        all_slides (Iterable[Slide]): An iterable containing all slides
        pattern (str): String to match

    Returns:
        SlideSubset: Subset of slides with slide number (1-indexed) as keys
    """
    subset = dict()
    for i, slide in enumerate(all_slides, 1):
        shapes = [*slide.shapes, *slide.slide_layout.shapes]  # type: ignore
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            if pattern in shape.text_frame.text or re.match(
                pattern, shape.text_frame.text
            ):
                subset[i] = slide
    return subset


def get_raw_text_extracts_from_slides(slides: SlideSubset) -> dict[int, list[str]]:
    """
    Returns the raw text from any shapes (including text boxes) in the provided slides.

    Args:
        slides (list[Slide]): List of slides

    Returns:
        dict[int, list[str]]: Raw string extracts according to slide number
    """
    result = {}
    for i, slide in slides.items():
        for shape in [*slide.shapes, *slide.slide_layout.shapes]:  # type: ignore
            if i in result:
                result[i].append(shape.text_frame.text)
            else:
                result[i] = [shape.text_frame.text]
    return result


class ContentChecker(BaseChecker):
    """
    Checks the content of the uploaded slides according to the inputs.
    """

    def __init__(
        self,
        file_path: str,
        presentation: Presentation,
        req_order_of_service: str,
        selected_date: str,
        sermon_discussion_qns: str,
    ) -> None:
        self.file_name = file_path
        self.presentation = presentation
        self.raw_req_order_of_service = req_order_of_service
        self.selected_date = selected_date
        self.sermon_discussion_qns = sermon_discussion_qns

    @cached_property
    def slides(self) -> list[Slide]:
        return [slide for slide in self.presentation.slides]  # type: ignore

    @cached_property
    def section_headers(self) -> SlideSubset:
        """
        Section header slides are slides that contain the text "order of service".
        This includes the Welcome slide, and all other item slides (i.e. Family Confession,
        Family Prayer etc.)

        Returns:
            SlideSubset: Subset of slides with slide number (1-indexed) as keys
        """
        text = "order of service"
        return get_slides_by_pattern(self.slides, text)

    @cached_property
    def sermon_discussion_slides(self) -> SlideSubset:
        """
        Sermon discussion slides are slides that contain the text "Sermon discussion questions".

        Returns:
            SlideSubset: Subset of slides with slide number (1-indexed) as keys
        """
        text = "Sermon discussion questions"
        return get_slides_by_pattern(self.slides, text)

    @cached_property
    def slide_order_of_service(self) -> SlideOrderOfService:
        """
        Returns the order of service for each slide as a dictionary where the keys are slide number,
            and the values are a list of strings.

        Returns:
            SlideOrderOfService: Order of service in each slide
        """

        def filter_order_of_service_only(
            texts: dict[int, list[str]]
        ) -> dict[int, list[str]]:
            result = {}
            for i, item_list in texts.items():
                for item in item_list:
                    if "order of service" in item:
                        result[i].append(item)
                    else:
                        result[i] = [item]
            return result

        def split_and_strip(text: str) -> list[str]:
            """
            Splits up raw text on newline characters and strips on both sides of the
            resulting string.
            """
            split_text = re.split("\n+", text)
            return [
                item.strip()
                for item in split_text
                if len(item) and "order of service" not in item
            ]

        section_header_text = get_raw_text_extracts_from_slides(self.section_headers)
        orders_of_service = filter_order_of_service_only(section_header_text)
        return {
            i: split_and_strip(item)
            for i, item_list in orders_of_service.items()
            for item in item_list
        }

    @cached_property
    def filtered_req_order_of_service(self) -> FilteredCleanOrderOfService:
        return filter_clean_req_order_of_service(
            get_clean_req_order_of_service(self.raw_req_order_of_service)
        )

    @cached_property
    def cleaned_sermon_discussion_qns(self) -> list[str]:
        return get_clean_sermon_discussion_qns(self.sermon_discussion_qns)

    def run(self) -> list[Result]:
        """
        Runs all the checks within the ContentChecker for a single Presentation instance.

        Returns:
            list[Result]: List of Result dictionaries
        """
        results = [
            self.check_existence_of_section_headers(),
            *self.check_section_headers_have_correct_order(),
            *self.check_all_dates_are_as_provided(),
            self.check_existence_of_lone_sermon_discussion_slide(),
            *self.check_sermon_discussion_qns_are_as_provided(),
        ]
        return self.sorted(results)

    def sorted(self, results: list[Result]) -> list[Result]:
        """
        Sorts results with the following precedence: Errors, Warnings, Passes.

        Args:
            results (list[Result]): Unsorted list of results

        Returns:
            list[Result]: Sorted list of results
        """
        return sorted(results, key=lambda x: x["status"], reverse=True)

    def check_existence_of_section_headers(self) -> Result:
        """
        Test that there exists at least 1 section header slide in the presentation.

        Returns:
            Result: Result of this test
        """
        result: Result = {
            "title": "Check existence of section header slides",
            "status": Status.PASS if len(self.section_headers) > 0 else Status.ERROR,
            "comments": f"Expected: >=1 section header slides. Provided: {len(self.section_headers)} section header slide(s) found.",
        }
        return result

    def check_section_headers_have_correct_order(
        self,
    ) -> list[Result]:
        """
        Test all section headers have the correct order of service by checking that the
        order of service provided in each section header slide matches the provided order
        of service.

        Section headers are the slides that separate each section of the service.
        These slides will usually have the order of service in grey on the right side.

        Returns:
            list[Result]: List of results for this test
        """
        # 1. Identify the section header slides
        # 2. Extract the text from the slide
        # 3. Check that the order of service from the text in the slides is correct
        results: list[Result] = []
        items_with_comments = (
            "(Opening Song|Closing Song|Hearing God(\u2018|\u2019|')s Word Read)"
        )

        for i, slide_text in self.slide_order_of_service.items():
            index = 0
            for entry in slide_text:

                title, comments = self.filtered_req_order_of_service[index]
                title = title.replace("\u2018", "\u2019")
                required_item = f"{title} \u2013 {comments}"
                is_commented_item = re.match(items_with_comments, entry)
                partial_ratio = fuzz.partial_ratio(required_item, entry)

                if is_commented_item is not None:
                    is_commented_items_correct = required_item == entry
                    if is_commented_items_correct:
                        index += 1
                        continue
                    elif 90 < partial_ratio < 100:
                        result = {
                            "title": "Check section headers are in the correct order: Is there a typo?",
                            "status": Status.WARNING,
                            "comments": f"On Slide {i}, Expected: '{required_item}'. Provided: '{entry}'. Similarity score = {partial_ratio} of 100",
                        }
                        index += 1
                    else:
                        result = {
                            "title": "Check section headers are in the correct order",
                            "status": Status.ERROR,
                            "comments": f"On Slide {i}, Expected: '{required_item}'. Provided: '{entry}'. Similarity score = {partial_ratio} of 100",
                        }
                    if result not in results:
                        results.append(result)
                elif "\u2018" in entry:
                    result = {
                        "title": "Check section headers are in the correct order: Is there a typo?",
                        "status": Status.WARNING,
                        "comments": f"On Slide {i}, Expected: '{required_item}'. Provided: '{entry}'. The use of the unicode character U+2018 (\u2018) is triggering this warning; replace this character with U+2019 (\u2019) or a standard single quote (') to resolve this error. Similarity score = {partial_ratio} of 100",
                    }
                    results.append(result)
                    index += 1
                elif entry.replace("\u2018", "\u2019") != title:
                    continue
                else:
                    index += 1

        if len(results) == 0:
            result: Result = {
                "title": "Check all required order of service items are present and in the correct order",
                "status": Status.PASS,
                "comments": "All slides containing order of service have the required order of service items and are presented in the correct order.",
            }
            results.append(result)
        return results

    def check_family_confession_content_matches_number(self) -> list[Result]:
        """
        Test family confession has the correct contents by checking that the words on the
        slide match the required content (usually a #number) specified in the order of
        service.

        Family Confession is an item in the order of service.

        Returns:
            list[Result]: List of Result dictionaries
        """
        # 1. Identify the family confession slides
        # 2. Check that there is a text box on the page that contains the family confession
        # 3. Check that it matches the number
        raise NotImplementedError

    def check_all_lyric_slides_have_no_title(self) -> list[Result]:
        """
        Test all lyric slides do not contain a title.

        Lyric slides are slides for the opening and closing song.

        Returns:
            list[Result]: List of Result dictionaries
        """
        # 1. Identify the lyric slides for the opening and closing song
        # 2. Check that the title is not present (i.e. none of the text boxes contain the song title exclusively)
        raise NotImplementedError

    def check_all_dates_are_as_provided(self) -> list[Result]:
        """
        Test all dates that appear in the presentation are equal to the date of Sunday service.

        Searches for certain date patterns in the slides.

        Date patterns:
            "01-Jan-2022"
            "01 Jan 2022"

        Returns:
            list[Result]: List of Result dictionaries
        """
        date_pattern = "\\d+[\\s-][A-Za-z]+[\\s-]\\d+"
        slides_with_dates = get_slides_by_pattern(self.slides, date_pattern)
        results = []
        for i, item_list in get_raw_text_extracts_from_slides(
            slides_with_dates
        ).items():
            for item in item_list:
                partial_ratio = fuzz.partial_ratio(item, self.selected_date)
                if re.match(date_pattern, item) and item.replace(
                    "_", " "
                ) != self.selected_date.replace("_", " "):
                    result = {
                        "title": "Check all dates that appear in the slides are the same as the date of Sunday service.",
                        "status": Status.ERROR,
                        "comments": f"On slide {i}, Expected: '{self.selected_date}'. Provided: '{item}'. Similarity score = {partial_ratio} of 100",
                    }
                    results.append(result)

        if len(results) == 0:
            result = {
                "title": "Check all dates that appear in the slides are the same as the date of Sunday service.",
                "status": Status.PASS,
                "comments": "All slides containing dates display the required date of Sunday service.",
            }
            results.append(result)

        return results

    def check_existence_of_lone_sermon_discussion_slide(self) -> Result:
        """
        Test that there exists exactly 1 sermon discussion slide in the presentation.

        Returns:
            Result: Result of this test
        """
        result: Result = {
            "title": "Check existence of sermon discussion slides.",
            "status": Status.PASS
            if len(self.sermon_discussion_slides) == 1
            else Status.ERROR,
            "comments": f"Expected: 1 sermon discussion slide. Provided: {len(self.sermon_discussion_slides)} sermon discussion slide(s) found.",
        }
        return result

    def check_sermon_discussion_qns_are_as_provided(self) -> list[Result]:
        """
        Test that the required questions appear in the sermon discussion slide.
        Ignores numbering (i.e. 1. and 2.) and checks sentences directly.

        Returns:
            list[Result]: List of Result dictionaries
        """
        # 1. Get the sermon discussion slides
        # 2. Check questions are in these slides
        raw_text_extracts = get_raw_text_extracts_from_slides(
            self.sermon_discussion_slides
        )
        slide_number = list(raw_text_extracts.keys())[0]
        raw_text = list(raw_text_extracts.values())[0]
        split_text = [item for item_list in raw_text for item in item_list.split("\n")]

        results = []
        for required_qn in self.cleaned_sermon_discussion_qns:
            for entry in split_text:
                partial_ratio = fuzz.partial_ratio(required_qn, entry)
                if required_qn in split_text:
                    break
                elif 90 < partial_ratio < 100:
                    result = {
                        "title": "Check sermon discussion questions are as provided: Is there a typo?",
                        "status": Status.WARNING,
                        "comments": f"On Slide {slide_number}, Expected: '{required_qn}'. Provided: '{entry}'. Similarity score = {partial_ratio} of 100",
                    }
                    results.append(result)
                    break
            # using for-else construct that runs if no break statement is hit
            else:
                # a good match for the required_qn was not found
                result: Result = {
                    "title": "Check sermon discussion questions are as provided.",
                    "status": Status.ERROR,
                    "comments": f"On Slide {slide_number}, Expected: '{required_qn}'. Could not find this required question.",
                }
                results.append(result)

        if len(results) == 0:
            result: Result = {
                "title": "Check sermon discussion questions are as provided.",
                "status": Status.PASS,
                "comments": "All sermon discussion questions are present.",
            }
            results.append(result)
        return results


class MultiContentChecker(BaseMultiChecker):
    """
    Extends ContentChecker for multiple presentation files.
    """

    def __init__(
        self,
        presentations: dict[str, Presentation],
        req_order_of_service: str,
        selected_date: str,
        sermon_discussion_qns: str,
    ) -> None:
        self.presentations = presentations
        self.req_order_of_service = req_order_of_service
        self.selected_date = selected_date
        self.sermon_discussion_qns = sermon_discussion_qns

    @cached_property
    def checkers(self) -> dict[str, ContentChecker]:
        return {
            file_name: ContentChecker(
                file_path=file_name,
                presentation=pptx,
                req_order_of_service=self.req_order_of_service,
                selected_date=self.selected_date,
                sermon_discussion_qns=self.sermon_discussion_qns,
            )
            for file_name, pptx in self.presentations.items()
        }

    def run(self) -> list[FileResults]:
        file_results = []
        for file_name, checker in self.checkers.items():
            file_results.append({"filename": file_name, "results": checker.run()})
        return file_results


if __name__ == "__main__":
    filename = "22.05 (10.30am) service slides.pptx"
    sermon_discussion_qns = """1. How have you been confronted with your own arrogance before God today? How have you been challenged to repent?
2. How has our passage been a comfort if we are seeking to live for God in this anti-God world?"""
    selected_date = "22 May 2022"

    with open(f"input/{filename}", "rb") as f:
        pptx = PresentationConstructor(f)

    slides: list[Slide] = [slide for slide in pptx.slides]  # type: ignore
    cc = ContentChecker(
        selected_date=selected_date,
        req_order_of_service=raw_req_order_of_service_no_declaration(),
        sermon_discussion_qns=sermon_discussion_qns,
        file_path=filename,
        presentation=pptx,
    )
    cc_results = cc.run()

    mcc = MultiContentChecker(
        presentations={filename: pptx},
        req_order_of_service=raw_req_order_of_service_no_declaration(),
        selected_date=selected_date,
        sermon_discussion_qns=sermon_discussion_qns,
    )
    mcc_results = mcc.run()
